import streamlit as st
import google.generativeai as genai
import os
import io
import base64
from PIL import Image
import pandas as pd
import matplotlib.pyplot as plt
import numpy as np
from datetime import datetime
import pdfkit
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
import tempfile
from docx import Document
import pdfplumber
from pptx import Presentation

# Configure Gemini API
GEMINI_API_KEY = "AIzaSyApb3uEAndgRMBIvIYGULdoO3xFcYcYJeU"
genai.configure(api_key=GEMINI_API_KEY)

# Set up page configuration
st.set_page_config(
    page_title="Exam Paper Handwriting Analyzer",
    page_icon="📝",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for better UI
st.markdown("""
<style>
    .main-header {
        font-size: 2.5rem;
        color: #1E3A8A;
        text-align: center;
        margin-bottom: 1rem;
        font-weight: bold;
    }
    .sub-header {
        font-size: 1.5rem;
        color: #2563EB;
        margin-bottom: 1rem;
    }
    .result-header {
        font-size: 1.8rem;
        color: #1E40AF;
        text-align: center;
        margin: 1rem 0;
        font-weight: bold;
    }
    .correct-answer {
        background-color: #DCFCE7;
        padding: 0.5rem;
        border-radius: 0.5rem;
        border-left: 4px solid #22C55E;
    }
    .incorrect-answer {
        background-color: #FEE2E2;
        padding: 0.5rem;
        border-radius: 0.5rem;
        border-left: 4px solid #EF4444;
    }
    .partial-answer {
        background-color: #FEF3C7;
        padding: 0.5rem;
        border-radius: 0.5rem;
        border-left: 4px solid #F59E0B;
    }
    .comment-box {
        background-color: #EFF6FF;
        padding: 0.75rem;
        border-radius: 0.5rem;
        border: 1px solid #BFDBFE;
        margin-top: 0.5rem;
    }
    .stTabs [data-baseweb="tab-list"] {
        gap: 2px;
    }
    .stTabs [data-baseweb="tab"] {
        background-color: white;
        border-radius: 4px 4px 0px 0px;
        border-right: 1px solid #E2E8F0;
        border-left: 1px solid #E2E8F0;
        border-top: 1px solid #E2E8F0;
        border-bottom: none;
        padding: 10px 16px;
        font-weight: 600;
    }
    .stTabs [aria-selected="true"] {
        background-color: #EFF6FF;
        border-bottom: 2px solid #2563EB;
    }
    .uploadedFile {
        border: 1px solid #E2E8F0;
        border-radius: 5px;
        padding: 10px;
        margin-bottom: 10px;
    }
    .stButton>button {
        background-color: #2563EB;
        color: white;
        font-weight: 600;
        border-radius: 4px;
        padding: 0.5rem 1rem;
        border: none;
    }
    .stButton>button:hover {
        background-color: #1E40AF;
    }
    div[data-testid="stExpander"] div[role="button"] p {
        font-weight: 600;
        color: #2563EB;
    }
</style>
""", unsafe_allow_html=True)

# Initialize session state variables
if 'questions_answers' not in st.session_state:
    st.session_state.questions_answers = []
if 'analyzed_results' not in st.session_state:
    st.session_state.analyzed_results = []
if 'total_score' not in st.session_state:
    st.session_state.total_score = 0
if 'max_score' not in st.session_state:
    st.session_state.max_score = 0
if 'student_info' not in st.session_state:
    st.session_state.student_info = {
        'name': '',
        'id': '',
        'subject': '',
        'exam_date': datetime.now().date(),
        'class': ''
    }
if 'evaluation_comments' not in st.session_state:
    st.session_state.evaluation_comments = ''

# Header section
st.markdown('<p class="main-header">Exam Paper Handwriting Analyzer</p>', unsafe_allow_html=True)

# Sidebar for navigation and app info
with st.sidebar:
    st.image("https://cdn-icons-png.flaticon.com/512/2631/2631448.png", width=100)
    st.title("Navigation")

    app_mode = st.radio(
        "Select Mode",
        ["Instructions", "Setup Exam", "Analyze Papers", "Results & Reports"]
    )

    st.divider()
    st.subheader("About")
    st.markdown("""
    **Exam Paper Analyzer** uses AI to evaluate handwritten exam papers, allocate marks, and generate comprehensive reports.

    **Features:**
    - Analyze handwritten answers
    - Auto-evaluate against model answers
    - Provide detailed feedback
    - Generate downloadable reports

    Made with ❤️ using Streamlit and Gemini AI
    """)

# Instructions page
if app_mode == "Instructions":
    st.markdown('<p class="sub-header">How to Use the Exam Paper Analyzer</p>', unsafe_allow_html=True)

    st.info("This application helps teachers evaluate handwritten exam papers using AI technology.")

    col1, col2 = st.columns(2)

    with col1:
        st.markdown("### Step 1: Setup Exam Details")
        st.markdown("""
        - Enter student information
        - Define question bank with model answers
        - Set marks for each question
        """)

        st.markdown("### Step 2: Upload & Analyze")
        st.markdown("""
        - Upload scanned exam papers
        - The AI will recognize handwriting
        - Compare with model answers
        """)

    with col2:
        st.markdown("### Step 3: Review & Adjust")
        st.markdown("""
        - Review AI evaluation
        - Adjust marks if needed
        - Add comments and feedback
        """)

        st.markdown("### Step 4: Generate Reports")
        st.markdown("""
        - View detailed results
        - Generate downloadable PDF reports
        - Save evaluation data
        """)

    st.warning("Ensure that uploaded images are clear for better handwriting recognition.")

    st.markdown("### Tips for Best Results")
    st.markdown("""
    - Use high-quality scans of handwritten papers
    - Provide detailed model answers for better comparison
    - Review AI evaluations for occasional adjustments
    - Use the comment feature to provide personalized feedback
    """)

# Setup Exam page
elif app_mode == "Setup Exam":
    st.markdown('<p class="sub-header">Setup Exam Parameters</p>', unsafe_allow_html=True)

    # Student information
    st.subheader("Student Information")
    col1, col2 = st.columns(2)
    with col1:
        st.session_state.student_info['name'] = st.text_input("Student Name", st.session_state.student_info['name'])
        st.session_state.student_info['subject'] = st.text_input("Subject", st.session_state.student_info['subject'])
        st.session_state.student_info['class'] = st.text_input("Class/Grade", st.session_state.student_info['class'])
    with col2:
        st.session_state.student_info['id'] = st.text_input("Student ID", st.session_state.student_info['id'])
        st.session_state.student_info['exam_date'] = st.date_input("Exam Date", st.session_state.student_info['exam_date'])

    # Question setup
    st.subheader("Question Bank Setup")
    st.markdown("Define questions, model answers, and mark allocation for the exam.")

    # Method to add a new question
    def add_question():
        st.session_state.questions_answers.append({
            'question': '',
            'model_answer': '',
            'max_marks': 0,
            'question_type': 'Short Answer'
        })

    # Method to remove a question
    def remove_question(idx):
        st.session_state.questions_answers.pop(idx)

    # Add first question if none exist
    if not st.session_state.questions_answers:
        add_question()

    # Display existing questions with edit capability
    for i, qa in enumerate(st.session_state.questions_answers):
        with st.expander(f"Question {i+1}", expanded=(i==len(st.session_state.questions_answers)-1)):
            col1, col2, col3 = st.columns([3, 1, 0.5])

            with col1:
                st.session_state.questions_answers[i]['question'] = st.text_area(
                    f"Question Text #{i+1}",
                    qa['question'],
                    height=100,
                    key=f"q_text_{i}"
                )

            with col2:
                st.session_state.questions_answers[i]['question_type'] = st.selectbox(
                    "Question Type",
                    ["Short Answer", "Long Answer", "Multiple Choice", "True/False"],
                    index=["Short Answer", "Long Answer", "Multiple Choice", "True/False"].index(qa['question_type']),
                    key=f"q_type_{i}"
                )

                st.session_state.questions_answers[i]['max_marks'] = st.number_input(
                    "Max Marks",
                    min_value=1,
                    max_value=100,
                    value=int(qa['max_marks']) if qa['max_marks'] > 0 else 5,
                    key=f"marks_{i}"
                )

            with col3:
                if len(st.session_state.questions_answers) > 1:
                    if st.button("Remove", key=f"remove_{i}"):
                        remove_question(i)
                        st.rerun()

            st.session_state.questions_answers[i]['model_answer'] = st.text_area(
                f"Model Answer #{i+1}",
                qa['model_answer'],
                height=150,
                key=f"model_{i}"
            )

    # Add button for new questions
    if st.button("+ Add Another Question"):
        add_question()
        st.rerun()

    # Calculate total possible marks
    st.session_state.max_score = sum(q['max_marks'] for q in st.session_state.questions_answers)
    st.info(f"Total exam marks: {st.session_state.max_score}")

    # Save setup
    if st.button("Save Exam Setup"):
        if not st.session_state.student_info['name']:
            st.error("Please enter student name")
        elif not st.session_state.student_info['subject']:
            st.error("Please enter subject name")
        elif not any(q['question'] and q['model_answer'] for q in st.session_state.questions_answers):
            st.error("Please enter at least one question and model answer")
        else:
            st.success("Exam setup saved successfully! Proceed to 'Analyze Papers' to continue.")

# Analyze Papers page
elif app_mode == "Analyze Papers":
    st.markdown('<p class="sub-header">Analyze Student Papers</p>', unsafe_allow_html=True)

    # Check if exam is set up
    if not st.session_state.questions_answers or not any(q['question'] for q in st.session_state.questions_answers):
        st.warning("Please set up the exam questions first in the 'Setup Exam' section.")
        st.stop()

    # Function to analyze handwriting with Gemini AI
    def analyze_handwriting(text, question, model_answer):
        # Initialize Gemini Vision model
        model = genai.GenerativeModel('gemini-1.5-flash')

        # Create the prompt for analysis
        prompt = f"""
        You are an expert exam grader analyzing student answers.

        QUESTION: {question}

        MODEL ANSWER: {model_answer}

        Analyze the answer and provide the following:
        1. Evaluation of the answer compared to the model answer (correct/partially correct/incorrect)
        2. Justification for your evaluation
        3. Suggested score as a percentage (0-100%)

        Format your response as JSON:
        {{
            "evaluation": "correct/partially correct/incorrect",
            "justification": "Your detailed justification",
            "score_percentage": 85,
            "feedback": "Constructive feedback for the student"
        }}
        """

        try:
            # Get response from Gemini
            response = model.generate_content([prompt, text])
            result = response.text

            # Extract JSON part if needed (sometimes Gemini wraps the JSON in markdown)
            if "```json" in result:
                result = result.split("```json")[1].split("```")[0].strip()

            # Convert to dictionary
            import json
            result_dict = json.loads(result)
            return result_dict
        except Exception as e:
            st.error(f"Error analyzing handwriting: {str(e)}")
            return {
                "evaluation": "error",
                "justification": f"Error: {str(e)}",
                "score_percentage": 0,
                "feedback": "Unable to analyze"
            }

    # Function to extract text from DOC files
    def extract_text_from_doc(file):
        doc = Document(file)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return '\n'.join(full_text)

    # Function to extract text from PDF files
    def extract_text_from_pdf(file):
        with pdfplumber.open(file) as pdf:
            full_text = []
            for page in pdf.pages:
                full_text.append(page.extract_text())
        return '\n'.join(full_text)

    # Function to extract text from PPT files
    def extract_text_from_ppt(file):
        prs = Presentation(file)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return '\n'.join(full_text)

    # Display current questions
    st.subheader("Current Exam Questions")
    for i, qa in enumerate(st.session_state.questions_answers):
        st.markdown(f"**Question {i+1}:** {qa['question']} *(Max marks: {qa['max_marks']})*")

    st.divider()

    # Upload and analyze section
    st.subheader("Upload Answer Sheets")

    uploaded_files = {}

    tabs = st.tabs([f"Question {i+1}" for i in range(len(st.session_state.questions_answers))])

    for i, tab in enumerate(tabs):
        with tab:
            st.write(f"**Question:** {st.session_state.questions_answers[i]['question']}")
            st.write(f"**Max Marks:** {st.session_state.questions_answers[i]['max_marks']}")

            uploaded_file = st.file_uploader(
                f"Upload answer for Question {i+1} (DOC, PDF, PPT)",
                type=["docx", "pdf", "pptx"],
                key=f"upload_{i}"
            )

            if uploaded_file is not None:
                file_extension = uploaded_file.name.split('.')[-1].lower()
                if file_extension == 'docx':
                    text = extract_text_from_doc(uploaded_file)
                elif file_extension == 'pdf':
                    text = extract_text_from_pdf(uploaded_file)
                elif file_extension == 'pptx':
                    text = extract_text_from_ppt(uploaded_file)
                else:
                    st.error("Unsupported file format")
                    continue

                st.text_area("Extracted Text", text, height=200)
                uploaded_files[i] = text

                # Analyze button for this question
                if st.button(f"Analyze Answer for Question {i+1}", key=f"analyze_{i}"):
                    with st.spinner("Analyzing handwriting..."):
                        # Get analysis from Gemini
                        result = analyze_handwriting(
                            text,
                            st.session_state.questions_answers[i]['question'],
                            st.session_state.questions_answers[i]['model_answer']
                        )

                        # Store results
                        if len(st.session_state.analyzed_results) <= i:
                            st.session_state.analyzed_results.extend([None] * (i + 1 - len(st.session_state.analyzed_results)))

                        # Calculate actual marks based on percentage
                        max_marks = st.session_state.questions_answers[i]['max_marks']
                        actual_marks = round((result['score_percentage'] / 100) * max_marks, 1)

                        result['actual_marks'] = actual_marks
                        result['max_marks'] = max_marks
                        result['question_index'] = i

                        st.session_state.analyzed_results[i] = result

                        # Display results
                        evaluation_class = "correct-answer" if result['evaluation'] == "correct" else \
                                         "partial-answer" if result['evaluation'] == "partially correct" else \
                                         "incorrect-answer"

                        st.markdown(f"<div class='{evaluation_class}'>", unsafe_allow_html=True)
                        st.markdown(f"**Evaluation:** {result['evaluation'].title()}")
                        st.markdown(f"**Score:** {actual_marks}/{max_marks} ({result['score_percentage']}%)")
                        st.markdown("</div>", unsafe_allow_html=True)

                        with st.expander("View Detailed Feedback"):
                            st.markdown(f"**Justification:** {result['justification']}")
                            st.markdown(f"**Feedback for Student:** {result['feedback']}")

                            # Allow teacher to adjust marks
                            adjusted_marks = st.slider(
                                "Adjust marks if needed:",
                                0.0,
                                float(max_marks),
                                float(actual_marks),
                                0.5,
                                key=f"adjust_{i}"
                            )

                            if adjusted_marks != actual_marks:
                                st.session_state.analyzed_results[i]['actual_marks'] = adjusted_marks
                                st.session_state.analyzed_results[i]['score_percentage'] = round((adjusted_marks / max_marks) * 100)
                                st.success(f"Marks adjusted to {adjusted_marks}/{max_marks}")
            else:
                st.info(f"Please upload the answer for Question {i+1}")

    # Analyze all button
    if uploaded_files and len(uploaded_files) == len(st.session_state.questions_answers):
        if st.button("Analyze All Remaining Questions"):
            progress_bar = st.progress(0)

            for i, qa in enumerate(st.session_state.questions_answers):
                if i not in uploaded_files or i >= len(st.session_state.analyzed_results) or st.session_state.analyzed_results[i] is None:
                    continue

                with st.spinner(f"Analyzing Question {i+1}..."):
                    # Get analysis from Gemini
                    result = analyze_handwriting(
                        uploaded_files[i],
                        qa['question'],
                        qa['model_answer']
                    )

                    # Store results
                    if len(st.session_state.analyzed_results) <= i:
                        st.session_state.analyzed_results.extend([None] * (i + 1 - len(st.session_state.analyzed_results)))

                    # Calculate actual marks based on percentage
                    max_marks = qa['max_marks']
                    actual_marks = round((result['score_percentage'] / 100) * max_marks, 1)

                    result['actual_marks'] = actual_marks
                    result['max_marks'] = max_marks
                    result['question_index'] = i

                    st.session_state.analyzed_results[i] = result

                # Update progress
                progress_bar.progress((i + 1)/len(st.session_state.questions_answers))

            st.success("All questions analyzed! Go to 'Results & Reports' to see the complete evaluation.")

    # If there are any analyzed results, show summary and option to go to results
    if any(st.session_state.analyzed_results):
        st.divider()
        st.subheader("Analysis Summary")

        # Calculate total score
        analyzed_count = sum(1 for r in st.session_state.analyzed_results if r is not None)
        total_score = sum(r['actual_marks'] for r in st.session_state.analyzed_results if r is not None)
        max_possible = sum(r['max_marks'] for r in st.session_state.analyzed_results if r is not None)

        st.session_state.total_score = total_score

        st.info(f"Analyzed {analyzed_count} out of {len(st.session_state.questions_answers)} questions")
        st.success(f"Current score: {total_score}/{max_possible} ({round((total_score/max_possible)*100, 1)}%)")

        # Option for teacher's overall comments
        st.session_state.evaluation_comments = st.text_area(
            "Overall Comments/Feedback",
            st.session_state.evaluation_comments,
            height=100
        )

        st.button("View Complete Results", on_click=lambda: st.switch_page(None))  # This doesn't actually switch pages, just visual

# Results & Reports page
elif app_mode == "Results & Reports":
    st.markdown('<p class="sub-header">Exam Results & Reports</p>', unsafe_allow_html=True)

    # Check if exam has been analyzed
    if not any(st.session_state.analyzed_results):
        st.warning("No exam has been analyzed yet. Please go to 'Analyze Papers' section first.")
        st.stop()

    # Student Info
    st.subheader("Student Information")
    col1, col2 = st.columns(2)
    with col1:
        st.write(f"**Name:** {st.session_state.student_info['name']}")
        st.write(f"**Subject:** {st.session_state.student_info['subject']}")
        st.write(f"**Class/Grade:** {st.session_state.student_info['class']}")
    with col2:
        st.write(f"**Student ID:** {st.session_state.student_info['id']}")
        st.write(f"**Exam Date:** {st.session_state.student_info['exam_date']}")

    st.divider()

    # Calculate total score and performance stats
    analyzed_questions = [r for r in st.session_state.analyzed_results if r is not None]
    total_obtained = sum(r['actual_marks'] for r in analyzed_questions)
    total_possible = sum(r['max_marks'] for r in analyzed_questions)
    percentage = (total_obtained / total_possible) * 100 if total_possible > 0 else 0

    # Determine grade based on percentage
    grade = "A+" if percentage >= 90 else \
            "A" if percentage >= 80 else \
            "B" if percentage >= 70 else \
            "C" if percentage >= 60 else \
            "D" if percentage >= 50 else "F"

    # Display score summary
    st.markdown('<p class="result-header">Exam Results Summary</p>', unsafe_allow_html=True)

    col1, col2, col3 = st.columns(3)
    with col1:
        st.metric("Total Score", f"{total_obtained}/{total_possible}")
    with col2:
        st.metric("Percentage", f"{percentage:.1f}%")
    with col3:
        st.metric("Grade", grade)

    # Detailed results
    st.subheader("Detailed Question Analysis")

    for i, result in enumerate(st.session_state.analyzed_results):
        if result is None:
            continue

        with st.expander(f"Question {i+1} - {result['actual_marks']}/{result['max_marks']} marks"):
            st.write(f"**Question:** {st.session_state.questions_answers[i]['question']}")
            st.write(f"**Model Answer:** {st.session_state.questions_answers[i]['model_answer']}")

            # Evaluation result with appropriate styling
            evaluation_class = "correct-answer" if result['evaluation'] == "correct" else \
                               "partial-answer" if result['evaluation'] == "partially correct" else \
                               "incorrect-answer"

            st.markdown(f"<div class='{evaluation_class}'>", unsafe_allow_html=True)
            st.write(f"**Evaluation:** {result['evaluation'].title()}")
            st.write(f"**Score:** {result['actual_marks']}/{result['max_marks']} ({result['score_percentage']}%)")
            st.markdown("</div>", unsafe_allow_html=True)

            st.markdown("<div class='comment-box'>", unsafe_allow_html=True)
            st.write("**Feedback:**")
            st.write(result['feedback'])
            st.markdown("</div>", unsafe_allow_html=True)

    # Performance visualization
    st.subheader("Performance Visualization")

    fig, ax = plt.subplots(figsize=(10, 6))

    # Question-wise score comparison
    question_nums = [f"Q{i+1}" for i, r in enumerate(st.session_state.analyzed_results) if r is not None]
    obtained_scores = [r['actual_marks'] for r in st.session_state.analyzed_results if r is not None]
    max_scores = [r['max_marks'] for r in st.session_state.analyzed_results if r is not None]
    percentages = [(o/m)*100 for o, m in zip(obtained_scores, max_scores)]

    x = np.arange(len(question_nums))
    width = 0.35

    ax.bar(x - width/2, obtained_scores, width, label='Obtained Marks')
    ax.bar(x + width/2, max_scores, width, label='Maximum Marks')

    for i, p in enumerate(percentages):
        ax.text(i, obtained_scores[i] + 0.1, f"{p:.0f}%", ha='center')

    ax.set_xticks(x)
    ax.set_xticklabels(question_nums)
    ax.set_ylabel('Marks')
    ax.set_title('Question-wise Performance')
    ax.legend()

    st.pyplot(fig)

    # Generate report functionality
    st.subheader("Generate Report")

    def create_pdf_report():
        # Create a temporary file to store the PDF
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_file:
            pdf_path = tmp_file.name

        # Set up the PDF document
        doc = SimpleDocTemplate(
            pdf_path,
            pagesize=letter,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=72
        )

        # Create styles
        styles = getSampleStyleSheet()
        title_style = styles['Title']
        heading_style = styles['Heading1']
        subheading_style = styles['Heading2']
        normal_style = styles['Normal']

        # Custom styles
        centered_style = ParagraphStyle(
            'centered',
            parent=styles['Normal'],
            alignment=1,
            spaceAfter=12
        )

        feedback_style = ParagraphStyle(
            'feedback',
            parent=styles['Normal'],
            leftIndent=20,
            spaceBefore=6,
            spaceAfter=6
        )

        # Build the document content
        elements = []

        # Title
        elements.append(Paragraph("Exam Evaluation Report", title_style))
        elements.append(Spacer(1, 0.25 * inch))

        # Student info
        student_info = [
            [Paragraph("<b>Student Name:</b>", normal_style), Paragraph(st.session_state.student_info['name'], normal_style)],
            [Paragraph("<b>Student ID:</b>", normal_style), Paragraph(st.session_state.student_info['id'], normal_style)],
            [Paragraph("<b>Class/Grade:</b>", normal_style), Paragraph(st.session_state.student_info['class'], normal_style)],
            [Paragraph("<b>Subject:</b>", normal_style), Paragraph(st.session_state.student_info['subject'], normal_style)],
            [Paragraph("<b>Exam Date:</b>", normal_style), Paragraph(str(st.session_state.student_info['exam_date']), normal_style)]
        ]

        info_table = Table(student_info, colWidths=[1.5*inch, 4*inch])
        info_table.setStyle(TableStyle([
            ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
            ('BACKGROUND', (0, 0), (0, -1), colors.lightgrey),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('PADDING', (0, 0), (-1, -1), 6)
        ]))

        elements.append(info_table)
        elements.append(Spacer(1, 0.25 * inch))

        # Results summary
        elements.append(Paragraph("Results Summary", heading_style))

        summary_data = [
            [Paragraph("<b>Total Score:</b>", normal_style),
             Paragraph(f"{total_obtained}/{total_possible}", normal_style)],
            [Paragraph("<b>Percentage:</b>", normal_style),
             Paragraph(f"{percentage:.1f}%", normal_style)],
            [Paragraph("<b>Grade:</b>", normal_style),
             Paragraph(grade, normal_style)]
        ]

        summary_table = Table(summary_data, colWidths=[1.5*inch, 4*inch])
        summary_table.setStyle(TableStyle([
            ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
            ('BACKGROUND', (0, 0), (0, -1), colors.lightgrey),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('PADDING', (0, 0), (-1, -1), 6)
        ]))

        elements.append(summary_table)
        elements.append(Spacer(1, 0.25 * inch))

        # Detailed question analysis
        elements.append(Paragraph("Detailed Question Analysis", heading_style))

        for i, result in enumerate(st.session_state.analyzed_results):
            if result is None:
                continue

            # Question header
            elements.append(Paragraph(f"Question {i+1}", subheading_style))

            # Question and model answer
            elements.append(Paragraph(f"<b>Question:</b> {st.session_state.questions_answers[i]['question']}", normal_style))
            elements.append(Paragraph(f"<b>Model Answer:</b> {st.session_state.questions_answers[i]['model_answer']}", normal_style))
            elements.append(Spacer(1, 0.1 * inch))

            # Student's answer
            elements.append(Paragraph(f"<b>Student's Answer (Transcribed):</b>", normal_style))
            elements.append(Paragraph(result['transcription'], feedback_style))
            elements.append(Spacer(1, 0.1 * inch))

            # Evaluation
            eval_color = colors.green if result['evaluation'] == "correct" else \
                         colors.orange if result['evaluation'] == "partially correct" else \
                         colors.red

            eval_data = [
                [Paragraph("<b>Evaluation:</b>", normal_style),
                 Paragraph(f"<font color='{eval_color}'>{result['evaluation'].title()}</font>", normal_style)],
                [Paragraph("<b>Score:</b>", normal_style),
                 Paragraph(f"{result['actual_marks']}/{result['max_marks']} ({result['score_percentage']}%)", normal_style)]
            ]

            eval_table = Table(eval_data, colWidths=[1.5*inch, 4*inch])
            eval_table.setStyle(TableStyle([
                ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                ('BACKGROUND', (0, 0), (0, -1), colors.lightgrey),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ('PADDING', (0, 0), (-1, -1), 6)
            ]))

            elements.append(eval_table)
            elements.append(Spacer(1, 0.1 * inch))

            # Feedback
            elements.append(Paragraph("<b>Feedback:</b>", normal_style))
            elements.append(Paragraph(result['feedback'], feedback_style))
            elements.append(Spacer(1, 0.25 * inch))

        # Overall comments
        if st.session_state.evaluation_comments:
            elements.append(Paragraph("Overall Comments", heading_style))
            elements.append(Paragraph(st.session_state.evaluation_comments, normal_style))

        # Build the PDF
        doc.build(elements)
        return pdf_path

    if st.button("Generate PDF Report"):
        with st.spinner("Generating PDF report..."):
            try:
                pdf_path = create_pdf_report()

                # Read the generated PDF file
                with open(pdf_path, "rb") as f:
                    pdf_bytes = f.read()

                # Create download button
                st.download_button(
                    label="Download PDF Report",
                    data=pdf_bytes,
                    file_name=f"{st.session_state.student_info['name']}_{st.session_state.student_info['subject']}_Report.pdf",
                    mime="application/pdf"
                )

                st.success("PDF report generated successfully!")
            except Exception as e:
                st.error(f"Error generating report: {str(e)}")

    # Export data option
    st.subheader("Export Data")

    def export_to_csv():
        # Create DataFrame for export
        data = []

        for i, result in enumerate(st.session_state.analyzed_results):
            if result is None:
                continue

            data.append({
                'Question Number': i+1,
                'Question Text': st.session_state.questions_answers[i]['question'],
                'Student Answer': result['transcription'],
                'Evaluation': result['evaluation'],
                'Score': result['actual_marks'],
                'Max Score': result['max_marks'],
                'Percentage': result['score_percentage'],
                'Feedback': result['feedback']
            })

        df = pd.DataFrame(data)

        # Add summary row
        summary = pd.DataFrame([{
            'Question Number': 'TOTAL',
            'Question Text': '',
            'Student Answer': '',
            'Evaluation': '',
            'Score': total_obtained,
            'Max Score': total_possible,
            'Percentage': percentage,
            'Feedback': st.session_state.evaluation_comments
        }])

        df = pd.concat([df, summary], ignore_index=True)

        # Convert to CSV
        csv = df.to_csv(index=False)
        return csv

    if st.button("Export to CSV"):
        try:
            csv = export_to_csv()

            st.download_button(
                label="Download CSV Data",
                data=csv,
                file_name=f"{st.session_state.student_info['name']}_{st.session_state.student_info['subject']}_Data.csv",
                mime="text/csv"
            )

            st.success("Data exported successfully!")
        except Exception as e:
            st.error(f"Error exporting data: {str(e)}")

    # Clear results option
    if st.button("Clear Results and Start New Analysis"):
        st.session_state.analyzed_results = []
        st.session_state.total_score = 0
        st.session_state.evaluation_comments = ''
        st.success("Results cleared. You can now start a new analysis.")
        st.rerun()
